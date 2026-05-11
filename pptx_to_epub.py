#!/usr/bin/env python3
"""
PowerPoint to Markdown converter.

Extracts text content from `.pptx` slides and builds structured Markdown
files that are easier for AI tools to ingest. Supports single-file and
whole-directory processing with both CLI and tkinter GUI entry points.
"""

from __future__ import annotations

import argparse
import statistics
import sys
import threading
from dataclasses import dataclass
from pathlib import Path
from queue import Empty, Queue

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER


@dataclass
class SlideBlock:
    kind: str
    text: str = ""
    level: int = 0
    rows: list[list[str]] | None = None


@dataclass
class SlideContent:
    title: str
    blocks: list[SlideBlock]


def clean_text(value: str) -> str:
    return " ".join((value or "").replace("\u00a0", " ").split())


def escape_cell(value: str) -> str:
    return clean_text(value).replace("|", "\\|")


def iter_text_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_text_shapes(shape.shapes)
            continue

        if getattr(shape, "has_table", False) or getattr(shape, "has_text_frame", False):
            yield shape


def placeholder_type(shape):
    if not getattr(shape, "is_placeholder", False):
        return None

    try:
        return shape.placeholder_format.type
    except Exception:
        return None


def is_title_placeholder(shape) -> bool:
    return placeholder_type(shape) in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}


def is_subtitle_placeholder(shape) -> bool:
    return placeholder_type(shape) == PP_PLACEHOLDER.SUBTITLE


def paragraph_text(paragraph) -> str:
    runs = "".join(run.text for run in paragraph.runs) if paragraph.runs else paragraph.text
    return clean_text(runs)


def paragraph_font_sizes(shape) -> list[float]:
    sizes: list[float] = []
    if not getattr(shape, "has_text_frame", False):
        return sizes

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            size = getattr(getattr(run, "font", None), "size", None)
            if size:
                try:
                    sizes.append(float(size.pt))
                except Exception:
                    pass
    return sizes


def paragraph_is_bold(paragraph) -> bool:
    bold_flags = [
        bool(run.font.bold)
        for run in paragraph.runs
        if getattr(getattr(run, "font", None), "bold", None) is not None
    ]
    return bool(bold_flags) and all(bold_flags)


def paragraph_is_heading(text: str, paragraph, base_font_size: float | None, shape) -> bool:
    if not text or len(text) > 80:
        return False
    if text.endswith((".", "!", "?", ";")):
        return False
    if getattr(paragraph, "level", 0):
        return False
    if is_subtitle_placeholder(shape):
        return True

    para_sizes = [
        float(run.font.size.pt)
        for run in paragraph.runs
        if getattr(getattr(run, "font", None), "size", None) is not None
    ]
    para_size = max(para_sizes) if para_sizes else None
    words = [word for word in text.split() if any(ch.isalpha() for ch in word)]
    title_like = words and sum(1 for word in words if word[:1].isupper()) >= max(1, len(words) // 2)

    if paragraph_is_bold(paragraph) and title_like:
        return True
    if para_size and base_font_size and para_size >= base_font_size * 1.2 and title_like:
        return True
    return False


def extract_shape_blocks(shape, base_font_size: float | None) -> list[SlideBlock]:
    blocks: list[SlideBlock] = []

    if getattr(shape, "has_table", False):
        rows: list[list[str]] = []
        for row in shape.table.rows:
            values = [escape_cell(cell.text) for cell in row.cells]
            if any(values):
                rows.append(values)
        if rows:
            blocks.append(SlideBlock(kind="table", rows=rows))
        return blocks

    if not getattr(shape, "has_text_frame", False):
        return blocks

    for paragraph in shape.text_frame.paragraphs:
        text = paragraph_text(paragraph)
        if not text:
            continue

        level = max(0, int(getattr(paragraph, "level", 0) or 0))
        if paragraph_is_heading(text, paragraph, base_font_size, shape):
            blocks.append(SlideBlock(kind="heading", text=text, level=3))
        elif level > 0:
            blocks.append(SlideBlock(kind="list", text=text, level=level))
        else:
            blocks.append(SlideBlock(kind="paragraph", text=text))
    return blocks


def presentation_title(presentation: Presentation, source_path: Path) -> str:
    core_title = clean_text(getattr(presentation.core_properties, "title", "") or "")
    if core_title:
        return core_title

    for index, slide in enumerate(presentation.slides, start=1):
        extracted = extract_slide_content(slide, index)
        if extracted.title and not extracted.title.startswith("Slide "):
            return extracted.title

    return source_path.stem.replace("_", " ").strip() or source_path.stem


def extract_slide_content(slide, slide_number: int) -> SlideContent:
    title = ""
    blocks: list[SlideBlock] = []

    all_sizes: list[float] = []
    for shape in iter_text_shapes(slide.shapes):
        all_sizes.extend(paragraph_font_sizes(shape))
    base_font_size = statistics.median(all_sizes) if all_sizes else None

    for shape in iter_text_shapes(slide.shapes):
        shape_blocks = extract_shape_blocks(shape, base_font_size)
        if not shape_blocks:
            continue

        if is_title_placeholder(shape) and not title:
            title = shape_blocks[0].text
            shape_blocks = shape_blocks[1:]

        blocks.extend(shape_blocks)

    if not title:
        title = f"Slide {slide_number}"

    return SlideContent(title=title, blocks=blocks)


def render_table(rows: list[list[str]]) -> list[str]:
    width = max(len(row) for row in rows)
    padded_rows = [row + [""] * (width - len(row)) for row in rows]
    header = padded_rows[0]
    separator = ["---"] * width
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    for row in padded_rows[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return lines


def render_slide_markdown(index: int, slide: SlideContent) -> list[str]:
    lines = [f"## {index:02d}. {slide.title}", ""]
    if not slide.blocks:
        lines.extend(["_No extractable text on this slide._", ""])
        return lines

    paragraph_buffer: list[str] = []

    def flush_paragraph():
        if paragraph_buffer:
            lines.append(" ".join(paragraph_buffer))
            lines.append("")
            paragraph_buffer.clear()

    for block in slide.blocks:
        if block.kind == "paragraph":
            paragraph_buffer.append(block.text)
            continue

        flush_paragraph()

        if block.kind == "heading":
            lines.extend([f"### {block.text}", ""])
        elif block.kind == "list":
            indent = "  " * max(0, block.level - 1)
            lines.append(f"{indent}- {block.text}")
        elif block.kind == "table" and block.rows:
            lines.extend(render_table(block.rows))
            lines.append("")

    flush_paragraph()
    if lines and lines[-1] != "":
        lines.append("")
    return lines


class PowerPointToMarkdownConverter:
    def extract_slides(self, pptx_path: Path) -> tuple[str, list[SlideContent]]:
        presentation = Presentation(str(pptx_path))
        title = presentation_title(presentation, pptx_path)
        slides = [
            extract_slide_content(slide, index)
            for index, slide in enumerate(presentation.slides, start=1)
        ]
        return title, slides

    def build_markdown(self, pptx_path: Path, output_path: Path) -> Path:
        title, slides = self.extract_slides(pptx_path)
        lines = [f"# {title}", ""]

        for index, slide in enumerate(slides, start=1):
            lines.extend(render_slide_markdown(index, slide))

        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text("\n".join(lines).rstrip() + "\n", encoding="utf-8")
        return output_path

    def collect_inputs(self, input_path: Path) -> list[Path]:
        if input_path.is_file():
            if input_path.suffix.lower() != ".pptx":
                raise ValueError("Only .pptx files are supported.")
            return [input_path]

        if input_path.is_dir():
            files = sorted(path for path in input_path.rglob("*.pptx") if path.is_file())
            if not files:
                raise ValueError(f"No .pptx files found in {input_path}")
            return files

        raise ValueError(f"Input path not found: {input_path}")

    def convert(self, input_path: Path, output_dir: Path, progress_callback=None) -> list[Path]:
        files = self.collect_inputs(input_path)
        results: list[Path] = []
        base_dir = input_path if input_path.is_dir() else input_path.parent

        if progress_callback:
            progress_callback(0, len(files), None, None, f"Found {len(files)} PowerPoint file(s).")

        for index, pptx_path in enumerate(files, start=1):
            if input_path.is_dir():
                relative_path = pptx_path.relative_to(base_dir).with_suffix(".md")
                output_path = output_dir / relative_path
            else:
                output_path = output_dir / f"{pptx_path.stem}.md"

            result = self.build_markdown(pptx_path, output_path)
            results.append(result)

            if progress_callback:
                progress_callback(index, len(files), pptx_path, result, None)

        return results


def launch_gui():
    import tkinter as tk
    from tkinter import filedialog, messagebox, scrolledtext, ttk

    class App:
        def __init__(self, root):
            self.root = root
            self.root.title("PowerPoint to Markdown")
            self.root.geometry("860x620")

            self.mode = tk.StringVar(value="file")
            self.input_path = tk.StringVar()
            self.output_dir = tk.StringVar()
            self.status = tk.StringVar(value="Ready")
            self.progress = tk.DoubleVar(value=0.0)
            self.queue = Queue()
            self.worker = None

            self._build()
            self.root.after(100, self._drain_queue)

        def _build(self):
            pad = dict(padx=8, pady=4)

            main = ttk.Frame(self.root, padding=10)
            main.grid(row=0, column=0, sticky="nsew")
            self.root.columnconfigure(0, weight=1)
            self.root.rowconfigure(0, weight=1)
            main.columnconfigure(1, weight=1)
            main.rowconfigure(5, weight=1)

            ttk.Label(main, text="PowerPoint to Markdown", font=("Segoe UI", 16, "bold")).grid(
                row=0, column=0, columnspan=3, sticky="w", **pad
            )
            ttk.Label(
                main,
                text="Extract slide structure into readable Markdown for AI and documentation workflows.",
                foreground="#555",
            ).grid(row=1, column=0, columnspan=3, sticky="w", **pad)

            mode_frame = ttk.LabelFrame(main, text="Input Type", padding=8)
            mode_frame.grid(row=2, column=0, columnspan=3, sticky="ew", **pad)
            ttk.Radiobutton(
                mode_frame,
                text="Single PowerPoint file",
                variable=self.mode,
                value="file",
                command=self._sync_defaults,
            ).grid(row=0, column=0, sticky="w", padx=4, pady=2)
            ttk.Radiobutton(
                mode_frame,
                text="Whole directory of PowerPoint files",
                variable=self.mode,
                value="directory",
                command=self._sync_defaults,
            ).grid(row=0, column=1, sticky="w", padx=12, pady=2)

            files = ttk.LabelFrame(main, text="Paths", padding=8)
            files.grid(row=3, column=0, columnspan=3, sticky="ew", **pad)
            files.columnconfigure(1, weight=1)

            self._file_row(files, 0, "PowerPoint file or folder:", self.input_path, self._browse_input)
            self._file_row(files, 1, "Output directory:", self.output_dir, self._browse_output)

            actions = ttk.Frame(main)
            actions.grid(row=4, column=0, columnspan=3, sticky="ew", **pad)
            self.convert_btn = ttk.Button(actions, text="Convert to Markdown", command=self._start)
            self.convert_btn.pack(side="left", padx=4)
            ttk.Button(actions, text="Clear Log", command=self._clear_log).pack(side="left", padx=4)

            ttk.Label(main, textvariable=self.status).grid(row=5, column=0, columnspan=3, sticky="w", **pad)
            ttk.Progressbar(main, maximum=100, variable=self.progress).grid(
                row=6, column=0, columnspan=3, sticky="ew", **pad
            )

            log_frame = ttk.LabelFrame(main, text="Activity Log", padding=8)
            log_frame.grid(row=7, column=0, columnspan=3, sticky="nsew", **pad)
            log_frame.columnconfigure(0, weight=1)
            log_frame.rowconfigure(0, weight=1)
            self.log = scrolledtext.ScrolledText(log_frame, height=20, font=("Consolas", 9))
            self.log.grid(row=0, column=0, sticky="nsew")

        def _file_row(self, parent, row, label, variable, command):
            ttk.Label(parent, text=label).grid(row=row, column=0, sticky="w", pady=2)
            ttk.Entry(parent, textvariable=variable).grid(row=row, column=1, sticky="ew", padx=4, pady=2)
            ttk.Button(parent, text="Browse", command=command).grid(row=row, column=2, padx=4, pady=2)

        def _browse_input(self):
            if self.mode.get() == "directory":
                selected = filedialog.askdirectory(title="Select PowerPoint Directory")
            else:
                selected = filedialog.askopenfilename(
                    title="Select PowerPoint File",
                    filetypes=[("PowerPoint files", "*.pptx")],
                )
            if selected:
                self.input_path.set(selected)
                self._sync_defaults()

        def _browse_output(self):
            selected = filedialog.askdirectory(title="Select Output Directory")
            if selected:
                self.output_dir.set(selected)

        def _sync_defaults(self):
            input_path = self.input_path.get().strip()
            if not input_path:
                return

            source = Path(input_path)
            if self.mode.get() == "directory":
                self.output_dir.set(str(source.parent / f"{source.name}_markdown"))
            else:
                self.output_dir.set(str(source.parent / "markdown"))

        def _append_log(self, message):
            self.log.insert("end", message + "\n")
            self.log.see("end")

        def _clear_log(self):
            self.log.delete("1.0", "end")

        def _start(self):
            input_value = self.input_path.get().strip()
            output_value = self.output_dir.get().strip()

            if not input_value:
                messagebox.showerror("Missing", "Please select a PowerPoint file or directory.")
                return
            if not output_value:
                messagebox.showerror("Missing", "Please select an output directory.")
                return
            if self.worker and self.worker.is_alive():
                return

            input_path = Path(input_value)
            output_dir = Path(output_value)

            self._clear_log()
            self.progress.set(0)
            self.status.set("Starting conversion...")
            self._append_log(f"Input: {input_path}")
            self._append_log(f"Output: {output_dir}")
            self.convert_btn.state(["disabled"])

            self.worker = threading.Thread(
                target=self._run_worker,
                args=(input_path, output_dir),
                daemon=True,
            )
            self.worker.start()

        def _run_worker(self, input_path: Path, output_dir: Path):
            converter = PowerPointToMarkdownConverter()

            def progress_callback(current, total, source_path, output_path, info_message):
                if info_message:
                    self.queue.put(("info", info_message))
                    return
                self.queue.put(
                    (
                        "progress",
                        (
                            current,
                            total,
                            f"Converted {source_path.name} -> {output_path.name}",
                        ),
                    )
                )

            try:
                results = converter.convert(input_path, output_dir, progress_callback=progress_callback)
                self.queue.put(("done", results))
            except Exception as exc:
                self.queue.put(("error", str(exc)))

        def _drain_queue(self):
            try:
                while True:
                    kind, payload = self.queue.get_nowait()
                    if kind == "info":
                        self._append_log(payload)
                        self.status.set(payload)
                    elif kind == "progress":
                        current, total, message = payload
                        percent = 0 if total <= 0 else (current / total) * 100
                        self.progress.set(percent)
                        self.status.set(message)
                        self._append_log(message)
                    elif kind == "done":
                        self.progress.set(100)
                        message = f"Finished. Created {len(payload)} Markdown file(s)."
                        self.status.set(message)
                        self._append_log(message)
                        self.convert_btn.state(["!disabled"])
                        messagebox.showinfo("Conversion Complete", message)
                    elif kind == "error":
                        self.status.set("Conversion failed.")
                        self._append_log(f"ERROR: {payload}")
                        self.convert_btn.state(["!disabled"])
                        messagebox.showerror("Conversion Failed", payload)
            except Empty:
                pass
            finally:
                self.root.after(100, self._drain_queue)

    root = tk.Tk()
    App(root)
    root.mainloop()


def build_arg_parser():
    parser = argparse.ArgumentParser(
        description="Convert PowerPoint slide text into structured Markdown files."
    )
    parser.add_argument("--gui", action="store_true", help="Launch the converter GUI.")
    parser.add_argument("--input", help="Path to a .pptx file or a directory containing .pptx files.")
    parser.add_argument("--output-dir", help="Directory where Markdown files will be written.")
    return parser


def run_cli(args) -> int:
    if not args.input:
        raise SystemExit("--input is required in CLI mode")
    if not args.output_dir:
        raise SystemExit("--output-dir is required in CLI mode")

    converter = PowerPointToMarkdownConverter()
    input_path = Path(args.input).expanduser()
    output_dir = Path(args.output_dir).expanduser()

    def progress_callback(current, total, source_path, output_path, info_message):
        if info_message:
            print(info_message)
            return
        print(f"[{current}/{total}] {source_path} -> {output_path}")

    results = converter.convert(input_path, output_dir, progress_callback=progress_callback)
    print(f"Created {len(results)} Markdown file(s) in {output_dir}")
    return 0


def main(argv=None) -> int:
    parser = build_arg_parser()
    argv = sys.argv[1:] if argv is None else argv

    if not argv or "--gui" in argv:
        launch_gui()
        return 0

    args = parser.parse_args(argv)
    return run_cli(args)


if __name__ == "__main__":
    raise SystemExit(main())
