"""
Smoke tests for the PowerPoint to Markdown converter.

Run with: python test_pptx_to_epub.py
"""

import sys
import tempfile
from pathlib import Path

from pptx import Presentation

from pptx_to_epub import PowerPointToMarkdownConverter


def create_sample_pptx(target: Path) -> Path:
    presentation = Presentation()

    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Introduction"
    body = slide.placeholders[1].text_frame
    body.text = "Overview paragraph"
    bullet = body.add_paragraph()
    bullet.text = "First bullet"
    bullet.level = 1

    slide2 = presentation.slides.add_slide(presentation.slide_layouts[5])
    textbox = slide2.shapes.add_textbox(left=1000000, top=1200000, width=6000000, height=3000000)
    frame = textbox.text_frame
    frame.text = "Body-only slide"
    frame.add_paragraph().text = "Follow-up point"

    presentation.save(target)
    return target


def assert_contains(text: str, expected: str):
    if expected not in text:
        raise AssertionError(f"Expected to find {expected!r}")


def test_single_file_conversion():
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        pptx_path = create_sample_pptx(temp_path / "deck.pptx")
        output_dir = temp_path / "out"

        converter = PowerPointToMarkdownConverter()
        results = converter.convert(pptx_path, output_dir)

        if len(results) != 1:
            raise AssertionError(f"Expected 1 Markdown file, got {len(results)}")

        markdown_path = results[0]
        if not markdown_path.exists():
            raise AssertionError("Markdown file was not created")

        content = markdown_path.read_text(encoding="utf-8")
        assert_contains(content, "# Introduction")
        assert_contains(content, "## 01. Introduction")
        assert_contains(content, "Overview paragraph")
        assert_contains(content, "- First bullet")


def main():
    tests = [("Single File Conversion", test_single_file_conversion)]
    failures = 0

    for name, func in tests:
        try:
            func()
            print(f"OK {name}")
        except Exception as exc:
            failures += 1
            print(f"FAIL {name}: {exc}")

    return 1 if failures else 0


if __name__ == "__main__":
    sys.exit(main())
